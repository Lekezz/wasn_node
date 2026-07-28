"""
make_presentation.py

Builds docs/summer-presentation.pptx and the figures it embeds, from the
content in docs/presentation-outline.md.

Why a script rather than a hand-built deck: the numbers on these slides come
from real measurements that are still changing (the 315 degree trials, the
board-versus-reference comparison once the firmware is flashed). Regenerating
is one command, so a slide never quietly keeps a stale number. The figures are
drawn here too, for the same reason.

Run:  python docs/make_presentation.py

WARNING: this OVERWRITES docs/summer-presentation.pptx. If you edit the deck in
PowerPoint, either stop running this script or copy your edits back into it.
Speaker notes are attached to each slide, so open the notes pane when
presenting.

Requires python-pptx and matplotlib.
"""

import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

BASE = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(BASE, "images")
OUT = os.path.join(BASE, "summer-presentation.pptx")

# Palette. Taken from the validated reference palette rather than picked by
# eye: slot 1 blue for data, slot 8 red reserved for the one status use (the
# flash limit that gets exceeded), neutral inks for everything textual.
INK = "#0b0b0b"
INK2 = "#52514e"
MUTED = "#a8a7a1"
ACCENT = "#2a78d6"
STATUS_BAD = "#e34948"
SURFACE = "#fcfcfb"
GRID = "#e3e2dd"

FONT = "Segoe UI"


def hexrgb(h):
    return RGBColor.from_string(h.lstrip("#").upper())


# ---------------------------------------------------------------- figures

def _style_axes(ax):
    """Recessive axes and grid, so the marks carry the chart."""
    ax.set_facecolor(SURFACE)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    for side in ("left", "bottom"):
        ax.spines[side].set_color(GRID)
    ax.tick_params(colors=INK2, labelsize=11, length=0)
    ax.grid(axis="y", color=GRID, linewidth=1, zorder=0)
    ax.set_axisbelow(True)


def _save(fig, name):
    path = os.path.join(IMG_DIR, name)
    fig.savefig(path, dpi=200, facecolor=SURFACE, bbox_inches="tight",
                pad_inches=0.25)
    plt.close(fig)
    print(f"  wrote {os.path.relpath(path, BASE)}")
    return path


def fig_tdoa():
    """The physics: a plane wave hits the four mics at different times."""
    fig, ax = plt.subplots(figsize=(7.0, 4.4))
    ax.set_facecolor(SURFACE)

    mics = np.array([[-4.6, 4.95], [4.6, 4.95], [-4.6, -4.95], [4.6, -4.95]])
    angle = np.deg2rad(35.0)
    direction = np.array([-np.cos(angle), -np.sin(angle)])   # travel direction

    # Wavefronts are lines perpendicular to the travel direction.
    perp = np.array([-direction[1], direction[0]])
    # Half-length 11 keeps the fronts clear of the caption strip at the
    # bottom; longer lines ran straight through the text.
    for k, offset in enumerate(np.arange(-1, 5) * 4.2):
        centre = -direction * offset + np.array([4.0, 4.0])
        p0 = centre - perp * 11
        p1 = centre + perp * 11
        ax.plot([p0[0], p1[0]], [p0[1], p1[1]], color=ACCENT,
                linewidth=2, alpha=0.45 if k == 0 else 0.22, zorder=1)

    # Arrival order is DERIVED from the direction, not asserted. Getting this
    # wrong by hand is easy and the slide would then teach the wrong thing:
    # the wave hits the most upstream mic first, so rank by projection onto
    # the direction the sound is coming FROM.
    upstream = -direction
    ranks = np.argsort(-(mics @ upstream))
    ordinal = {0: "1st", 1: "2nd", 2: "3rd", 3: "4th"}

    for rank, idx in enumerate(ranks):
        x, y = mics[idx]
        ax.scatter([x], [y], s=210, color=INK, zorder=4)
        ax.text(x, y, str(idx), color="white", fontsize=11.5, ha="center",
                va="center", zorder=5, fontweight="bold")
        # Top row labels go above, bottom row below, so nothing lands on a
        # wavefront label or on the caption.
        above = y > 0
        ax.text(x, y + (2.0 if above else -2.0),
                f"arrives {ordinal[rank]}", color=INK2, fontsize=11,
                ha="center", va="bottom" if above else "top",
                fontweight="bold" if rank == 0 else "normal")

    ax.annotate("", xy=(9.0, 5.6), xytext=(13.0, 8.4),
                arrowprops=dict(arrowstyle="-|>", color=ACCENT, linewidth=2.6))
    ax.text(13.4, 9.0, "incoming sound", color=ACCENT, fontsize=12,
            ha="right", va="bottom", fontweight="bold")

    ax.set_xlim(-11.5, 14.5)
    ax.set_ylim(-13.5, 11.5)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("Direction shows up as the ORDER and SPACING of arrivals",
                 color=INK, fontsize=13.5, pad=16, loc="left")
    ax.text(-11.5, -13.2,
            "1 sample at 16 kHz = 62.5 us = 2.1 cm of travel,\n"
            "so delays must be measured to a fraction of a sample",
            color=INK2, fontsize=11.5, ha="left", va="bottom")
    return _save(fig, "fig_tdoa.png")


def fig_array():
    """The built array, to scale, with the bearing convention."""
    fig, ax = plt.subplots(figsize=(6.4, 4.6))
    ax.set_facecolor(SURFACE)

    mics = {0: (-4.625, 4.95), 1: (4.625, 4.95),
            2: (-4.625, -4.95), 3: (4.625, -4.95)}

    for idx, (x, y) in mics.items():
        ax.scatter([x], [y], s=260, color=ACCENT, zorder=3)
        ax.text(x, y, str(idx), color="white", fontsize=12, ha="center",
                va="center", zorder=4, fontweight="bold")

    ax.plot([-4.625, 4.625, 4.625, -4.625, -4.625],
            [4.95, 4.95, -4.95, -4.95, 4.95],
            color=MUTED, linewidth=1.4, linestyle=(0, (5, 4)), zorder=1)

    ax.annotate("", xy=(4.625, 6.9), xytext=(-4.625, 6.9),
                arrowprops=dict(arrowstyle="<|-|>", color=INK2, linewidth=1.4))
    ax.text(0, 7.3, "9.25 cm", color=INK2, fontsize=11, ha="center")
    # Height dimension on the LEFT: the right side belongs to the +x bearing
    # arrow, and the two collided when both were on the right.
    ax.annotate("", xy=(-6.9, 4.95), xytext=(-6.9, -4.95),
                arrowprops=dict(arrowstyle="<|-|>", color=INK2, linewidth=1.4))
    ax.text(-7.4, 0, "9.9 cm", color=INK2, fontsize=11, va="center",
            ha="center", rotation=90)

    ax.annotate("", xy=(3.1, 0), xytext=(0, 0),
                arrowprops=dict(arrowstyle="-|>", color=INK, linewidth=2))
    ax.text(3.4, 0, "+x, 0 deg", color=INK, fontsize=11, va="center")
    ax.annotate("", xy=(0, 3.1), xytext=(0, 0),
                arrowprops=dict(arrowstyle="-|>", color=INK, linewidth=2))
    ax.text(0.35, 3.4, "90 deg", color=INK, fontsize=11)

    ax.set_xlim(-9.2, 9.4)
    ax.set_ylim(-8.5, 8.6)
    ax.set_aspect("equal")
    ax.axis("off")
    ax.set_title("Built array, measured port to port with calipers",
                 color=INK, fontsize=13, pad=10, loc="left")
    ax.text(-8, -8.2,
            "aperture 13.5 cm (6.32 samples)   condition number 1.07, "
            "so no blind directions",
            color=INK2, fontsize=10.5, ha="left", va="bottom")
    return _save(fig, "fig_array.png")


def fig_sync():
    """Sync proof: measured delay is linear in mic position."""
    delays = np.array([0.000, -0.481, -1.354, -1.801])
    pos = np.arange(4)

    fit = np.polyfit(pos, delays, 1)
    line = np.polyval(fit, pos)
    resid = np.max(np.abs(delays - line))

    fig, ax = plt.subplots(figsize=(6.6, 4.0))
    _style_axes(ax)

    xs = np.linspace(-0.25, 3.25, 50)
    ax.plot(xs, np.polyval(fit, xs), color=MUTED, linewidth=2,
            linestyle=(0, (5, 4)), zorder=2)
    ax.scatter(pos, delays, s=120, color=ACCENT, zorder=3)

    for x, y in zip(pos, delays):
        ax.annotate(f"{y:+.3f}", (x, y), textcoords="offset points",
                    xytext=(0, 13), ha="center", color=INK, fontsize=11)

    ax.text(2.55, -0.35, "straight-line fit", color=INK2, fontsize=11)
    ax.set_xticks(pos)
    ax.set_xticklabels([f"mic {i}" for i in pos])
    ax.set_ylabel("delay vs mic0 (samples)", color=INK2, fontsize=11)
    ax.set_ylim(-2.4, 0.75)
    ax.set_title("Channels start together: delay is linear in mic position",
                 color=INK, fontsize=13, pad=14, loc="left")
    ax.text(0, -2.28, f"worst departure from the line: {resid:.3f} samples",
            color=INK2, fontsize=11)
    return _save(fig, "fig_sync.png")


def fig_wall():
    """The headline result: the wall was the error, not the code."""
    labels = ["65 cm from wall", "2.1 m from wall"]
    errors = [3.40, 0.58]
    resid = [0.875, 0.119]
    colors = [MUTED, ACCENT]

    fig, axes = plt.subplots(1, 2, figsize=(8.2, 3.9))
    for ax, vals, title, unit in (
            (axes[0], errors, "Bearing error at a true 0 deg", "deg"),
            (axes[1], resid, "Worst triangle residual", "samples")):
        _style_axes(ax)
        bars = ax.bar(labels, vals, width=0.5, color=colors, zorder=3)
        for bar, v in zip(bars, vals):
            ax.annotate(f"{v:.3f} {unit}" if unit == "samples"
                        else f"{v:.2f} {unit}",
                        (bar.get_x() + bar.get_width() / 2, v),
                        textcoords="offset points", xytext=(0, 7),
                        ha="center", color=INK, fontsize=11.5,
                        fontweight="bold")
        ax.set_title(title, color=INK, fontsize=12.5, pad=12, loc="left")
        ax.set_ylim(0, max(vals) * 1.32)
        ax.tick_params(axis="x", labelsize=11)

    fig.suptitle("Same code, same array, same room. Only the wall moved.",
                 color=INK, fontsize=13.5, x=0.02, ha="left", y=1.04)
    return _save(fig, "fig_wall.png")


def fig_flash():
    """The CMSIS-DSP size problem and what fixed it."""
    labels = ["FFT tables\nas shipped", "FFT tables\nselective mode",
              "Whole DSP set\nas shipped", "Whole DSP set\ntrimmed",
              "Final\nimage"]
    kb = [707, 40, 712, 44, 55.7]
    colors = [MUTED, ACCENT, MUTED, ACCENT, ACCENT]

    fig, ax = plt.subplots(figsize=(8.0, 4.2))
    _style_axes(ax)
    bars = ax.bar(labels, kb, width=0.58, color=colors, zorder=3)
    for bar, v in zip(bars, kb):
        ax.annotate(f"{v:g} KB", (bar.get_x() + bar.get_width() / 2, v),
                    textcoords="offset points", xytext=(0, 7), ha="center",
                    color=INK, fontsize=11.5, fontweight="bold")

    ax.axhline(512, color=STATUS_BAD, linewidth=2, zorder=4)
    ax.text(4.45, 528, "512 KB, the whole flash on this part",
            color=STATUS_BAD, fontsize=11, ha="right", fontweight="bold")

    ax.set_ylabel("flash (KB)", color=INK2, fontsize=11)
    ax.set_ylim(0, 800)
    ax.tick_params(axis="x", labelsize=10.5)
    ax.set_title("The FFT tables alone did not fit in flash",
                 color=INK, fontsize=13, pad=14, loc="left")
    return _save(fig, "fig_flash.png")


def build_figures():
    os.makedirs(IMG_DIR, exist_ok=True)
    print("figures:")
    return {
        "tdoa": fig_tdoa(),
        "array": fig_array(),
        "sync": fig_sync(),
        "wall": fig_wall(),
        "flash": fig_flash(),
    }


# ------------------------------------------------------------------ deck

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_background(slide, hexcolor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hexrgb(hexcolor)


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _para(frame, text, size, color, bold=False, space_after=10,
          first=False, bullet=False, align=PP_ALIGN.LEFT):
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hexrgb(color)
    run.font.name = FONT
    return p


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _accent_bar(slide, color=ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62),
                                 Inches(1.28), Inches(1.05), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hexrgb(color)
    bar.line.fill.background()
    bar.shadow.inherit = False


def new_slide(prs, title=None, accent=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])     # blank
    _set_background(slide, SURFACE)
    if title:
        frame = _textbox(slide, Inches(0.62), Inches(0.45), Inches(12.1),
                         Inches(0.9))
        _para(frame, title, 30, INK, bold=True, first=True, space_after=0)
        if accent:
            _accent_bar(slide)
    return slide


def bullets_slide(prs, title, bullets, notes, width=12.1, left=0.62,
                  top=1.75, size=19):
    slide = new_slide(prs, title)
    frame = _textbox(slide, Inches(left), Inches(top), Inches(width),
                     Inches(4.9))
    for i, b in enumerate(bullets):
        indent = b.startswith("    ")
        _para(frame, b.strip(), size - (2 if indent else 0),
              INK2 if indent else INK, first=(i == 0), bullet=not indent,
              space_after=14 if not indent else 8)
    _notes(slide, notes)
    return slide


def image_slide(prs, title, image, bullets, notes, img_top=1.7,
                img_height=4.5):
    """Image on the left, short bullets on the right."""
    slide = new_slide(prs, title)
    pic = slide.shapes.add_picture(image, Inches(0.62), Inches(img_top),
                                   height=Inches(img_height))
    # Keep the picture inside the left two thirds; scale down if it is wide.
    max_w = Inches(7.9)
    if pic.width > max_w:
        ratio = max_w / pic.width
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)

    text_left = Emu(pic.left + pic.width) + Inches(0.35)
    avail = SLIDE_W - text_left - Inches(0.5)
    if bullets and avail > Inches(2.4):
        frame = _textbox(slide, text_left, Inches(img_top), avail,
                         Inches(4.8))
        for i, b in enumerate(bullets):
            _para(frame, b, 15.5, INK, first=(i == 0), bullet=True,
                  space_after=13)
    _notes(slide, notes)
    return slide


def full_image_slide(prs, title, image, caption, notes):
    slide = new_slide(prs, title)
    pic = slide.shapes.add_picture(image, Inches(0.62), Inches(1.65),
                                   height=Inches(4.7))
    if pic.width > Inches(12.1):
        ratio = Inches(12.1) / pic.width
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int((SLIDE_W - pic.width) / 2)
    if caption:
        frame = _textbox(slide, Inches(0.62), Inches(6.5), Inches(12.1),
                         Inches(0.7))
        _para(frame, caption, 15, INK2, first=True)
    _notes(slide, notes)
    return slide


def table_slide(prs, title, headers, rows, notes, caption=None,
                col_widths=None, highlight_row=None):
    slide = new_slide(prs, title)
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.42 * nrows)
    shape = slide.shapes.add_table(nrows, ncols, Inches(0.62), Inches(1.75),
                                   Inches(12.1), height)
    table = shape.table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(Inches(12.1) * w / total)

    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.name = FONT
        run.font.color.rgb = hexrgb("#ffffff")
        cell.fill.solid()
        cell.fill.fore_color.rgb = hexrgb(INK)

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(14)
            run.font.name = FONT
            emphasise = (highlight_row is not None and r - 1 == highlight_row)
            run.font.bold = emphasise
            run.font.color.rgb = hexrgb(ACCENT if emphasise else INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hexrgb(
                "#f2f1ec" if r % 2 else SURFACE)

    if caption:
        frame = _textbox(slide, Inches(0.62),
                         Inches(1.75) + height + Inches(0.25),
                         Inches(12.1), Inches(1.4))
        _para(frame, caption, 15, INK2, first=True)
    _notes(slide, notes)
    return slide


def build_deck(figs):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1 title
    slide = new_slide(prs, accent=False)
    frame = _textbox(slide, Inches(0.9), Inches(2.1), Inches(11.5),
                     Inches(2.6))
    _para(frame, "Acoustic Source Localization", 44, INK, bold=True,
          first=True, space_after=4)
    _para(frame, "on a Wireless Sensor Node", 44, ACCENT, bold=True,
          space_after=26)
    _para(frame, "Leke,  undergraduate EE,  University of Virginia", 20, INK2,
          space_after=6)
    _para(frame, "Supervisor: Ben", 20, INK2, space_after=6)
    _para(frame, "Summer 2026   |   work in progress", 18, ACCENT, bold=True,
          space_after=0)
    _notes(slide,
           "Set expectations in one sentence: this is a progress report on a "
           "project that works end to end but is not finished, and you will "
           "be clear about which parts are proven and which are not.\n\n"
           "If you have the array with you, hold it up here. It is small, and "
           "people are surprised by that.")

    # --- 2 goal
    bullets_slide(
        prs, "The goal",
        ["A small node that hears a sound and reports the direction it "
         "came from",
         "Several nodes cross their bearings to locate a source",
         "This summer: ONE node, microphone to bearing in degrees",
         "The network part is not started. That is honest, not a spoiler"],
        "Explain the eventual system briefly so the audience knows where this "
        "is heading, then narrow immediately to what was actually built. Say "
        "the network part is not started now, so it is not a surprise later.\n\n"
        "Reference design is SpeechCompass (CHI 2025), same microcontroller "
        "family and the same microphones.")

    # --- 3 physics
    image_slide(
        prs, "How direction comes out of timing", figs["tdoa"],
        ["Sound reaches the four mics at slightly different times",
         "Those differences depend on direction",
         "At 16 kHz one sample is 2.1 cm of travel",
         "So delays need sub-sample accuracy",
         "GCC-PHAT plus peak interpolation gets it"],
        "This is the physical core of the talk and worth spending time on. "
        "The punchline is sub-sample accuracy: whole-sample resolution is not "
        "enough, which is why the method is GCC-PHAT with an interpolated "
        "peak rather than simple peak finding.\n\n"
        "PHAT means the correlation is weighted by phase, which makes it "
        "robust to the very uneven frequency content of a clap. If someone "
        "asks why not plain cross correlation, that is the answer, and there "
        "is a backup slide.")

    # --- 4 system overview
    bullets_slide(
        prs, "System overview",
        ["Four PDM MEMS mics, 16 kHz per channel",
         "STM32L552 at 96 MHz: DFSDM decimation, DMA into RAM",
         "Clap trigger, one second stored",
         "Bearing computed ON the board with CMSIS-DSP",
         "Raw audio also dumped over serial, on purpose",
         "    Two outputs from one capture is what makes validation possible: "
         "the board's answer and the samples that produced it"],
        "Walk the chain left to right once: mics, DFSDM, DMA, RAM, clap "
        "trigger, then it splits into the on-board localizer and the raw dump "
        "to a PC.\n\n"
        "Emphasise that both paths exist deliberately. The board computes its "
        "own answer AND the raw audio comes out, so the same samples can be "
        "re-run against the reference implementation on a PC. Without that "
        "there would be no way to tell a firmware bug from a room problem.\n\n"
        "Consider drawing the block diagram on the whiteboard as you say it.")

    # --- 5 the SEL trick
    bullets_slide(
        prs, "Four microphones on two data pins",
        ["The chip has four DFSDM channels, not eight",
         "Each mic's SEL pin picks the rising or the falling clock edge",
         "Two mics with opposite SEL share one wire without colliding",
         "Mics 0 and 1 on PE7, mics 2 and 3 on PB1, one clock net to all four",
         "    This is the manufacturer's intended use of SEL, not a hack"],
        "A good 'clever bit' slide. The original plan assumed eight DFSDM "
        "channels; this part has four. Rather than change parts, the SEL pin "
        "fits four microphones into the pins available.\n\n"
        "Mic drives the line while the clock is high if SEL is grounded, and "
        "while it is low if SEL is at 3V3, so the two never collide. Channel "
        "pin redirection inside DFSDM is what lets both halves of each pair "
        "be reached.\n\n"
        "Finding the four-channel limit early rather than late is the reason "
        "the design fits at all.")

    # --- 6 the sync problem
    bullets_slide(
        prs, "The problem that had to be solved first",
        ["The four digital filters must start on the same clock edge",
         "If they do not, every channel carries an unknown offset",
         "An unknown offset is indistinguishable from real acoustic delay",
         "Fix: synchronized start, followers armed first, trigger filter last",
         "    The data cannot reveal this on its own, so it was tested "
         "directly"],
        "This is the single point of failure of the whole approach. "
        "Everything downstream measures time differences, so a constant "
        "per-channel offset poisons every result and still looks perfectly "
        "plausible.\n\n"
        "Make the point that the system cannot detect this from a normal "
        "recording. That is exactly why the next slide exists: a test "
        "designed so that only a sync fault could produce a certain result.")

    # --- 7 sync proof
    image_slide(
        prs, "Proving synchronization before trusting anything",
        figs["sync"],
        ["Mics packed in a tight line, about 2 cm apart",
         "True delay across the whole array is then under 3 samples",
         "So a large offset could ONLY be a sync fault",
         "Measured: linear in mic position to 0.131 samples",
         "Fixture retired the day it passed"],
        "The design of the test is the point worth making, more than the "
        "numbers. A spread out array would have confounded a sync fault with "
        "real geometry. Packing the mics tight made the two separable, "
        "because at 2 cm spacing physics cannot produce a large delay.\n\n"
        "The four measured delays were 0.000, -0.481, -1.354 and -1.801 "
        "samples, monotonic across the line and straight to within 0.131 "
        "samples.\n\n"
        "That fixture could never localize anything, its aperture is under "
        "3 samples, so it was retired immediately after it did its one job.")

    # --- 8 geometry
    image_slide(
        prs, "Array geometry, and why it is not hardcoded", figs["array"],
        ["Built array 9.25 x 9.9 cm, calipered port to port",
         "Layout changed twice mid-build, for physical reasons only",
         "So geometry lives in ONE registry file",
         "Condition number 1.07, so no blind directions",
         "Trying another layout is a one line change"],
        "The layout changed twice in a day: first only one breadboard was "
        "available, which forced a thin rectangle, then two were glued "
        "together and a near-square fit again. Nothing about the mathematics "
        "changed. That is why geometry was moved out of the analysis scripts "
        "into a registry with one active layout.\n\n"
        "Measure port to port, not pin to pin: on a breakout board the mic "
        "port is offset from the header pins, and 1 mm is about 0.047 samples.")

    # --- 9 geometry table
    table_slide(
        prs, "Mean error alone is misleading",
        ["Layout", "Aperture", "Condition number", "Mean", "p90", "Worst"],
        [["9.25 x 9.9 cm  (BUILT)", "13.5 cm", "1.07", "1.03 deg",
          "2.09 deg", "2.39 deg"],
         ["9.25 x 12 cm", "15.2 cm", "1.30", "0.70 deg", "2.01 deg",
          "2.16 deg"],
         ["9.25 x 16 cm", "18.5 cm", "1.73", "0.51 deg", "1.10 deg",
          "1.42 deg"],
         ["10 cm square (sim reference)", "14.1 cm", "1.00", "1.12 deg",
          "2.32 deg", "2.43 deg"],
         ["2.8 x 10 cm (one breadboard)", "10.4 cm", "3.57", "1.87 deg",
          "7.87 deg", "8.43 deg"],
         ["2.5 cm square", "3.5 cm", "1.00", "6.33 deg", "8.60 deg",
          "9.20 deg"]],
        "The row to point at is the 2.8 x 10 cm single-board rectangle. Its "
        "MEAN is a respectable 1.87 degrees, but its p90 is 7.87 degrees, "
        "because the error is not spread evenly: it piles up in four narrow "
        "cones about 30 degrees off the long axis. Averaged over all "
        "directions that looks fine. If your source sits in one of those "
        "cones it is not fine.\n\n"
        "Condition number predicts that directly, which is why the code warns "
        "based on it rather than on a hardcoded list of bad angles. Near 1 "
        "means uniform accuracy; above about 2 means blind directions exist.\n\n"
        "The extreme case: a straight line of mics is rank 1, the across-axis "
        "component is unconstrained, and the estimator returns 0 or 180 "
        "degrees regardless of truth. The script hard-fails rather than print "
        "a plausible number.",
        caption="Simulation only: 36 angles x 3 trials, 2 m source, 20 dB SNR. "
                "No reverberation, and the estimator is handed the exact "
                "geometry, so this RANKS layouts rather than predicts real "
                "accuracy.",
        col_widths=[3.4, 1.3, 1.9, 1.2, 1.2, 1.2],
        highlight_row=0)

    # --- 10 method
    bullets_slide(
        prs, "How correctness was established",
        ["Build ONE trusted reference, check everything against it",
         "1.  Verify in simulation, where truth is known",
         "2.  Prove channel sync separately from geometry",
         "3.  Share code between the test and the reference so they cannot "
         "drift",
         "4.  Compare the embedded port to the reference on identical samples",
         "    Simulation: delay bias under 0.001 samples, angle error under "
         "0.1 deg, about 0.4 deg at 0 dB SNR"],
        "This is the methodology slide and probably the most important one "
        "for a research audience.\n\n"
        "The temptation in a localization project is to judge a result by "
        "whether the angle looks about right. That does not work, because "
        "plenty of bugs produce plausible angles. A wrong channel map, a "
        "sign error, a filter offset: all of them give you a number in "
        "degrees that could be true.\n\n"
        "Point 3 is a small thing with a big payoff: the bring-up script "
        "imports GCC-PHAT from the reference rather than reimplementing it, "
        "so the test and the thing being tested cannot quietly drift apart.")

    # --- 11 result
    full_image_slide(
        prs, "Result, and what dominated the error", figs["wall"],
        "At 65 cm the wall reflection arrived inside the analysis window and "
        "corrupted one specific mic pair. Moving to 2.1 m fixed it with no "
        "code change at all.",
        "The best story in the project, so tell it as a story.\n\n"
        "The first ground-truthed capture was 3.4 degrees off and one mic "
        "pair was inconsistent with the other five. The array was 65 cm from "
        "a wall. The reflection was arriving soon enough to land inside the "
        "2048 sample analysis window, so the estimator was fitting a mixture "
        "of the direct sound and its echo.\n\n"
        "Moving the array to 2.1 m of clearance took the error to +0.58 "
        "degrees, which is at the simulation's own accuracy floor. Nothing in "
        "the code changed.\n\n"
        "This is why the capture folders record the room setup alongside the "
        "data: distance to the nearest wall is a property of the "
        "measurement, not of the run.")

    # --- 12 self check
    bullets_slide(
        prs, "A self-check that needs no ground truth",
        ["Delay is antisymmetric: A to B, plus B to C, must equal A to C",
         "It does not, when a reflection has corrupted a pair",
         "Worst violation across all four triangles flags a bad capture",
         "Residual tracked the error: 0.119, 0.347, 0.499 samples",
         "    against errors of 0.6, 2.1 and 6.4 degrees",
         "The firmware flagged both bad captures on its own"],
        "The key property is that this needs no knowledge of the true angle, "
        "which makes it usable in the field and not only in a validation "
        "run. In a deployed node there is no protractor.\n\n"
        "The firmware prints 'inconsistent: suspect a reflection' when the "
        "residual passes 0.3 samples, and it did so unprompted on both of "
        "the 315 degree trials before anyone had worked out why they were "
        "worse.\n\n"
        "It is a consistency check, not a correctness check: a capture can "
        "be self-consistent and still wrong. But an inconsistent one is "
        "definitely suspect.")

    # --- 13 embedding
    full_image_slide(
        prs, "Putting it on the microcontroller", figs["flash"],
        "Selective table mode names only the three tables a 4096 point real "
        "FFT needs. Final image 55.7 KB, with 48 KB of FFT buffers in a "
        "separate RAM section that costs no flash.",
        "Good slide for showing that embedded work is largely about budgets.\n\n"
        "CMSIS-DSP was vendored in by hand rather than added through the "
        "config tool, because the only available pack carrying it also "
        "switches on a USB device stack, a USB core and an AI application, "
        "all competing for the RAM the capture buffers need.\n\n"
        "Then the size problem: compiled whole, the tables file is 707 KB, "
        "larger than the entire 512 KB flash. It holds twiddle and "
        "bit-reversal tables for every FFT length and data type, and the "
        "generic initializer references all of them, so the linker cannot "
        "drop any. Naming only the three tables actually needed took it to "
        "40 KB.\n\n"
        "Timing estimate is about 10 ms of FFT work per 32 ms frame at "
        "96 MHz, so roughly 3x headroom.")

    # --- 14 silent failures
    table_slide(
        prs, "Every failure this summer was silent",
        ["What broke", "What you saw", "Actual cause"],
        [["Config tool regeneration", "Capture quietly wrong, twice",
          "DMA reverted to Normal mode, later to Byte width"],
         ["Serial dump", "Only part of the recording arrived",
          "Size argument is 16-bit: 128000 truncated to 62464"],
         ["Board's report", "Every float blank, integers fine",
          "nano.specs printf drops floating point support"],
         ["Clap trigger", "LED on forever, script timed out",
          "Threshold sized next to a wall, too high once it was gone"],
         ["Weak transient warning", "Loud clap called marginal",
          "DC offset in the noise estimate, not real room noise"]],
        "None of these announced themselves. Not one produced an error "
        "message. That is the honest engineering slide, and it is worth "
        "saying that this is what most of the debugging time went to.\n\n"
        "The blank float one is a nice puzzle to pose to the audience: every "
        "integer printed correctly and every float printed nothing. That "
        "exact pattern is the fingerprint of the small printf implementation "
        "linked with nano.specs, which omits float support unless you ask for "
        "it. Fixed by formatting floats with integer arithmetic instead, so "
        "it cannot come back from a build setting.\n\n"
        "What worked as defence: commit before and after every regeneration "
        "so a diff exposes it, treat every compiler warning as guilty, and "
        "build self-checks that flag a bad result without needing to know "
        "the right answer.",
        col_widths=[2.4, 3.4, 6.3])

    # --- 15 not done
    slide = new_slide(prs, "What is not done")
    frame = _textbox(slide, Inches(0.62), Inches(1.8), Inches(12.1),
                     Inches(4.9))
    items = [
        ("The board's own bearing has NEVER printed",
         "fix is written and committed, but not flashed yet"),
        ("Ground truth covers two angles, three trials",
         "a real validation plot needs a full sweep"),
        ("315 degrees is worse than 0 degrees and unexplained",
         "the array should have no bad directions"),
        ("Nothing is wireless, and there is one node",
         "multi-node needs inter-node time sync, a hard problem itself"),
        ("Claps only, one second at a time",
         "no continuous operation, no speech, no battery"),
    ]
    for i, (head, sub) in enumerate(items):
        _para(frame, head, 20, INK, bold=True, first=(i == 0), bullet=True,
              space_after=2)
        _para(frame, sub, 16, INK2, space_after=15)
    _notes(slide,
           "Do not rush this slide and do not apologise through it. Being "
           "precise about what is unproven is what makes the proven parts "
           "credible.\n\n"
           "The first item is the important one. The integer stages of the "
           "embedded port match the reference exactly, onset and analysis "
           "window were identical on all three captures, so find_clap is "
           "validated. But the arithmetic downstream of it has never produced "
           "a visible number on hardware. So describe the port as partly "
           "validated, never as working.\n\n"
           "If asked when it will be done: the fix is committed, it needs a "
           "flash and one capture, and the comparison script is already "
           "written and tested.")

    # --- 16 next steps
    bullets_slide(
        prs, "Next steps",
        ["1.  Flash the current firmware, confirm the report prints",
         "2.  Run the board versus reference comparison on a fresh capture",
         "3.  Sweep eight angles, five trials each, produce the validation "
         "plot",
         "4.  Then choose a direction:",
         "    multi-node cross-bearing  |  continuous operation  |  sources "
         "harder than claps"],
        "Steps 1 and 2 block everything else, and step 2 is already written "
        "and tested against synthetic data, so it is ready the moment the "
        "firmware is flashed.\n\n"
        "Close by naming the decision you actually want input on: which of "
        "the three directions in step 4 is most useful to the group. "
        "Multi-node is the most interesting but it needs time "
        "synchronization between nodes, which is a substantial problem in "
        "its own right and worth being explicit about rather than assuming.\n\n"
        "Then stop and take questions.")

    # --- backup divider
    slide = new_slide(prs, accent=False)
    frame = _textbox(slide, Inches(0.9), Inches(3.1), Inches(11.5),
                     Inches(1.4))
    _para(frame, "Backup", 40, MUTED, bold=True, first=True, space_after=8)
    _para(frame, "for questions", 22, MUTED, space_after=0)
    _notes(slide, "Everything after this point is for questions, not for the "
                  "main flow.")

    bullets_slide(
        prs, "Backup: why GCC-PHAT and not plain cross correlation",
        ["Plain cross correlation is dominated by whatever frequencies are "
         "loudest",
         "A clap is broadband but very uneven, and a room colours it further",
         "PHAT divides out the magnitude and keeps only phase",
         "That gives a sharp correlation peak instead of a broad one",
         "Peak is then interpolated for sub-sample resolution",
         "    Cost: PHAT is noisier at low SNR, since it trusts quiet "
         "frequencies as much as loud ones"],
        "Only if asked. The honest trade is in the last line: PHAT whitens "
        "everything, so at poor SNR it is weighting noise-dominated bins as "
        "heavily as clean ones. For a loud clap in a quiet room that is a "
        "good trade, which is the case this project targets.")

    bullets_slide(
        prs, "Backup: the signal chain in detail",
        ["PDM mic at 2.4 MHz, one bit per clock",
         "DFSDM sinc3 filter, oversampling ratio 150, gives 16 kHz",
         "Each 32-bit filter result is shifted right by 8 to make int16 PCM",
         "DMA in circular ping-pong halves, so capture never stops to copy",
         "All four channels drained as ONE aligned frame",
         "    That last point is what protects channel alignment, and "
         "GCC-PHAT depends on it entirely"],
        "The frame point is the subtle one. The code acts only when all four "
        "channels have flagged the same buffer half, then treats all four "
        "identically. If channels were handled independently, a clap "
        "arriving at the wrong moment could start one channel a frame ahead "
        "of another, and the estimator would read that bookkeeping error as "
        "acoustic delay.\n\n"
        "It is also why the warm-up discard is one global counter rather "
        "than four per-channel counters.")

    bullets_slide(
        prs, "Backup: why the array is planar, and what that costs",
        ["Four mics in a plane give BEARING only, not elevation",
         "A source above the plane is indistinguishable from its mirror below",
         "Elevation would need a non-planar layout, for example a tetrahedron",
         "For a node on a ceiling or a wall, bearing in the plane may be "
         "enough",
         "    This is the one limitation that is expensive to undo later"],
        "Worth raising with the group rather than defaulting into it. If the "
        "eventual deployment wants elevation, the array shape has to change, "
        "and every geometry number and the mounting change with it. Cheap to "
        "decide now, expensive to retrofit.")

    bullets_slide(
        prs, "Backup: how captures are stored",
        ["captures / session / angle / trial",
         "A session is one sitting with one room setup",
         "Session notes record the port, the geometry and the wall distance",
         "The board's own report is saved beside the samples that produced it",
         "    That is what makes board versus reference comparison possible "
         "after the fact"],
        "The room setup is grouped by session because it is a property of "
        "the data, not of the run. The wall distance determined whether a "
        "whole session was trustworthy, so it lives with the captures.\n\n"
        "Saving the board's report next to the samples was a late addition "
        "and it mattered: without it the report scrolled past in the "
        "terminal and was gone, because only one process can hold the serial "
        "port so you cannot watch it in a terminal at the same time.")

    prs.save(OUT)
    print(f"\ndeck: {os.path.relpath(OUT, BASE)}  "
          f"({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    figs = build_figures()
    build_deck(figs)
